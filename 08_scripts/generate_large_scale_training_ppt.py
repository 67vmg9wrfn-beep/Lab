#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

OUT_PPTX = "/Users/codex/Lab/07_reports/ppt/2024-01-29_large-scale-training-complete.pptx"
OUT_MD = "/Users/codex/Lab/07_reports/ppt/2024-01-29_large-scale-training-complete.slides.md"

# Palette (medical / signal processing style)
NAVY = RGBColor(14, 36, 64)
SLATE = RGBColor(42, 67, 101)
TEAL = RGBColor(0, 137, 123)
MINT = RGBColor(0, 191, 166)
LIGHT_BG = RGBColor(246, 250, 252)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(26, 32, 44)
GRAY = RGBColor(90, 102, 118)
ACCENT = RGBColor(255, 140, 66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, subtitle=None, dark=False):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(1.0))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Cambria"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE if dark else NAVY
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.7), Inches(1.38), Inches(12.0), Inches(0.55)).text_frame
        s.clear()
        sp = s.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Calibri"
        sp.font.size = Pt(16)
        sp.font.color.rgb = RGBColor(222, 230, 240) if dark else GRAY


def add_bullets(slide, x, y, w, h, items, color=DARK, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_card(slide, x, y, w, h, title, body, fill=WHITE):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = RGBColor(220, 230, 240)
    rect.line.width = Pt(1)

    t = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.4), Inches(0.4)).text_frame
    t.clear()
    tp = t.paragraphs[0]
    tp.text = title
    tp.font.name = "Calibri"
    tp.font.size = Pt(17)
    tp.font.bold = True
    tp.font.color.rgb = NAVY

    b = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.58), Inches(w - 0.4), Inches(h - 0.72)).text_frame
    b.clear()
    for i, line in enumerate(body):
        p = b.paragraphs[0] if i == 0 else b.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(4)


def add_footer(slide, text):
    f = slide.shapes.add_textbox(Inches(0.7), Inches(7.08), Inches(12.0), Inches(0.3)).text_frame
    f.clear()
    p = f.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY


def add_notes(slide, lines):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = "\n".join(lines)


slides_md = []

def md_slide(title, bullets):
    slides_md.append(f"## {title}\n")
    for b in bullets:
        slides_md.append(f"- {b}")
    slides_md.append("")

# 1 Title
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_title(s, "Large-Scale Training of Foundation Models for Wearable Biosignals", "ICLR 2024 · Apple · Focus: PPG/ECG AI pipeline for device-grade monitoring", dark=True)
add_bullets(s, 0.75, 2.2, 8.6, 2.8, [
    "From sparse labels -> self-supervised pretraining on longitudinal wearable data",
    "141,207 participants, ~20M PPG segments, ~3.75M ECG segments",
    "Reusable framework for algorithm architecture, training, and validation",
], color=RGBColor(232, 240, 248), size=21)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.8), Inches(1.8), Inches(2.8), Inches(3.7))
accent.fill.solid(); accent.fill.fore_color.rgb = TEAL; accent.line.fill.background()
kpi = s.shapes.add_textbox(Inches(10.1), Inches(2.05), Inches(2.2), Inches(3.2)).text_frame
for i, txt in enumerate(["141K", "Participants", "20M", "PPG Segments", "3.75M", "ECG Segments"]):
    p = kpi.paragraphs[0] if i == 0 else kpi.add_paragraph()
    p.text = txt
    p.font.name = "Calibri"
    p.font.bold = (i % 2 == 0)
    p.font.size = Pt(30 if i % 2 == 0 else 14)
    p.font.color.rgb = WHITE
add_footer(s, "Source: Abbaspourazad et al., ICLR 2024")
add_notes(s, ["开场：强调这不是单任务模型，而是平台级训练框架。", "重点数据规模和可迁移性。"])
md_slide("Slide 1 - Title", ["ICLR 2024 Apple biosignal foundation model framework", "Scale: 141K participants, 20M PPG, 3.75M ECG"])

# 2 Why now
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_title(s, "Why This Matters for Monitoring Device Teams")
add_card(s, 0.7, 1.6, 4.0, 2.2, "Engineering Problem", ["Clinical labels are expensive", "Signal quality is variable in free-living data", "Task-by-task model development does not scale"])
add_card(s, 4.95, 1.6, 4.0, 2.2, "Paper's Answer", ["Pretrain on unlabeled longitudinal biosignals", "Transfer embeddings to many downstream tasks", "Reduce dependence on dense labels"])
add_card(s, 9.2, 1.6, 3.4, 2.2, "Why You Care", ["Applicable to ECG/PPG/BP/SpO2 pipelines", "Supports architecture reuse", "Improves time-to-validation"])
add_bullets(s, 0.9, 4.2, 11.8, 2.4, [
    "Core thesis: representation quality is a first-class product artifact",
    "Practical impact: one pretrained encoder can service multiple biomarker heads",
], size=20)
add_footer(s, "Framework lens: data -> SSL pretraining -> downstream validation -> deployment")
add_notes(s, ["把这篇文章映射到监护设备研发流程。"])
md_slide("Slide 2 - Why it matters", ["Label scarcity and free-living noise block traditional supervised pipelines", "Foundation encoders improve reuse and speed"])

# 3 Data scale
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_title(s, "Dataset and Study Design (AHMS)")
chart_data = CategoryChartData()
chart_data.categories = ["Participants", "Segments (M)", "Avg days/participant", "Total span (days)"]
chart_data.add_series("PPG", [141.207, 19.854, 92.54, 890])
chart_data.add_series("ECG", [106.643, 3.744, 23.27, 1240])
chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.8), Inches(7.8), Inches(4.8), chart_data).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.value_axis.has_major_gridlines = True
add_card(s, 8.9, 1.9, 3.9, 4.5, "Data Facts", [
    "PPG: 60s segments @ 64/256Hz, 4 optical channels",
    "ECG: 30s lead-I strips @ 512Hz -> 128Hz",
    "Split: 80/10/10 by participants (no overlap)",
    "Informed consent + ongoing longitudinal study",
])
add_footer(s, "Table 1 from the paper")
add_notes(s, ["强调参与者级划分是验证可信性的关键。"])
md_slide("Slide 3 - Dataset", ["AHMS large-scale longitudinal wearable dataset", "Participant-level split and modality-specific preprocessing"])

# 4 pipeline
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_title(s, "Pretraining Pipeline (System View)")
steps = [
    (0.9, "1) Segment Curation", "Balance per participant\nMinimum segments constraint"),
    (3.4, "2) Stochastic Augmentation", "Cutout / warp / noise\nDifferent strengths for PPG vs ECG"),
    (5.9, "3) Encoder + Projection", "1D EfficientNet backbone\n256-d embedding -> 128-d projection"),
    (8.4, "4) Regularized InfoNCE", "Contrastive + KoLeo entropy reg\nL2-normalized embeddings"),
    (10.9, "5) Momentum Update", "EMA target branch\nStabilize representation learning"),
]
for x, t, b in steps:
    add_card(s, x, 2.1, 2.2, 3.6, t, b.split('\n'), fill=WHITE)
for x in [3.0, 5.5, 8.0, 10.5]:
    arr = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(3.55), Inches(0.3), Inches(0.5))
    arr.fill.solid(); arr.fill.fore_color.rgb = TEAL; arr.line.fill.background()
add_footer(s, "Algorithmic flow reconstructed from Sections 3-4")
add_notes(s, ["这里对应你们的算法架构图：数据层、表征层、任务层。"])
md_slide("Slide 4 - Pipeline", ["5-stage SSL pipeline with participant-level positive pairs", "Regularized InfoNCE + momentum branch"])

# 5 objective and augmentations
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_title(s, "Objective Function and Augmentation Policy")
add_card(s, 0.8, 1.7, 6.2, 2.3, "Loss Design", [
    "L = InfoNCE + KoLeo regularization",
    "Temperature = 0.04, KoLeo weight = 0.1",
    "32x A100 GPUs, Adam optimizer",
])
add_card(s, 0.8, 4.2, 6.2, 2.4, "Why this matters", [
    "InfoNCE aligns positive pairs and separates negatives",
    "KoLeo increases embedding entropy / spread",
    "Momentum branch improves stability for ECG-like low-variability signals",
])
add_card(s, 7.3, 1.7, 5.2, 4.9, "Augmentation probabilities", [
    "PPG: cutout 0.4, mag-warp 0.25, noise 0.25,",
    "      channel permute 0.25, time-warp 0.15",
    "ECG: cutout 0.8, mag-warp 0.5, noise 0.5,",
    "      time-warp 0.3 (no channel permute)",
    "Interpretation: ECG gets stronger distortion because",
    "within-participant variability is lower than PPG",
])
add_footer(s, "Section 3 + implementation details")
add_notes(s, ["这页可直接给算法工程师做训练配置模板。"])
md_slide("Slide 5 - Objective", ["Regularized InfoNCE objective and practical hyperparameters", "Modality-specific augmentation policy"])

# 6 key downstream results table2
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_title(s, "Downstream Results: Participant-level Positive Pairs Win")
add_card(s, 0.8, 1.7, 12.0, 1.0, "Table 2 headline", ["Switching from segment-level to participant-level pairs materially improves PPG and ECG downstream metrics"])
add_card(s, 0.8, 2.9, 5.8, 3.5, "PPG (participant-level)", [
    "Age cls AUC 0.976 (vs 0.900 segment-level)",
    "Age reg MAE 3.19 (vs 6.60)",
    "BMI cls AUC 0.918 (vs 0.766)",
    "BMI reg MAE 2.54 (vs 4.05)",
    "Sex cls AUC 0.993 (vs 0.870)",
])
add_card(s, 7.0, 2.9, 5.8, 3.5, "ECG (participant-level)", [
    "Age cls AUC 0.916 (vs 0.783 segment-level)",
    "Age reg MAE 6.33 (vs 9.01)",
    "BMI cls AUC 0.797 (vs 0.734)",
    "BMI reg MAE 3.72 (vs 4.11)",
    "Sex cls AUC 0.951 (vs 0.854)",
])
add_footer(s, "Direct numbers from Table 2")
add_notes(s, ["核心结论：正样本对选择策略是高杠杆设计点。"])
md_slide("Slide 6 - Results", ["Participant-level positive pairs outperform segment-level across PPG/ECG", "Large gains in age/BMI/sex probes"])

# 7 PPG vs ECG insights
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_title(s, "PPG vs ECG: Representation Behavior Differences")
add_card(s, 0.8, 1.8, 4.1, 4.7, "Observed in paper", [
    "ECG clusters are denser in t-SNE",
    "ECG InfoNCE loss is lower (easier pretraining)",
    "ECG dispersion ratio is smaller",
    "Yet PPG embeddings better predict many health targets",
])
add_card(s, 5.2, 1.8, 3.8, 4.7, "Interpretation", [
    "Signal modality matters",
    "A single SSL recipe may not be globally optimal",
    "Need modality-aware augmentation and objectives",
])
add_card(s, 9.3, 1.8, 3.9, 4.7, "Device-team implication", [
    "Do not assume one encoder policy fits all parameters",
    "Set separate training policies for PPG/ECG/BP/SpO2",
    "Track modality-specific validation dashboards",
])
add_footer(s, "Section 5.2.2 discussion")
add_notes(s, ["这页是方法论：多参数监护系统要分模态策略。"])
md_slide("Slide 7 - Modality distinction", ["ECG appears easier to contrastively pretrain, but PPG was more predictive for many targets", "Use modality-specific policies"])

# 8 ablation framework table3
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_title(s, "Ablation: Pretraining Framework Comparison")
chart_data = CategoryChartData()
chart_data.categories = ["Ours", "Ours(no KoLeo)", "SimCLR(var)", "BYOL(var)"]
chart_data.add_series("PPG SER", [104.13, 101.29, 98.87, 51.18])
chart_data.add_series("ECG SER", [113.82, 108.84, 103.65, 63.67])
chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.8), Inches(8.0), Inches(4.9), chart_data).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.value_axis.has_major_gridlines = True
add_card(s, 9.1, 1.8, 3.8, 4.9, "Takeaways", [
    "Their framework > SimCLR/BYOL variations",
    "KoLeo contributes measurable gains",
    "Contrastive approaches were more robust here",
    "Unsupervised proxy metrics are helpful but not sufficient",
])
add_footer(s, "Table 3 (smooth effective rank)")
add_notes(s, ["说明：SER是proxy，不等价于临床终点表现。"])
md_slide("Slide 8 - Ablation framework", ["Ours outperforms SimCLR/BYOL variations on SER", "KoLeo regularization is important"])

# 9 ablation encoder
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_title(s, "Ablation: Encoder Architecture vs Footprint")
add_card(s, 0.8, 1.8, 4.0, 4.8, "1D-EfficientNet", [
    "PPG SER 104.13 / ECG SER 113.82",
    "Params: 3.3M (PPG), 2.5M (ECG)",
    "Best efficiency-performance tradeoff",
])
add_card(s, 4.95, 1.8, 4.0, 4.8, "1D-ResNet", [
    "PPG SER 95.81 / ECG SER 111.92",
    "Params: 16.9M",
    "Larger but not clearly better",
])
add_card(s, 9.1, 1.8, 3.9, 4.8, "1D-ViT", [
    "PPG SER 104.89 / ECG SER 114.62",
    "Params: 7.2M",
    "Good performance, higher compute/memory",
])
add_footer(s, "Table 4")
add_notes(s, ["对设备端部署：先做EfficientNet系，再评估ViT收益是否值得。"])
md_slide("Slide 9 - Architecture tradeoff", ["EfficientNet-style 1D encoder gives strong efficiency", "ViT competitive but heavier"])

# 10 verification framework
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_title(s, "How to Reuse This as a Monitoring Algorithm Verification Framework")
add_card(s, 0.8, 1.8, 3.0, 4.9, "Layer 1\nRepresentation QA", [
    "SER / embedding drift",
    "Modality-specific pretraining loss",
    "Within-vs-across subject dispersion",
])
add_card(s, 4.0, 1.8, 3.0, 4.9, "Layer 2\nTask QA", [
    "Linear probes for key biomarkers",
    "Primary metrics + calibration",
    "Subgroup metrics by age/sex/BMI",
])
add_card(s, 7.2, 1.8, 3.0, 4.9, "Layer 3\nRobustness QA", [
    "Distribution shift slices",
    "Missingness / label shift checks",
    "Sensor-quality stress tests",
])
add_card(s, 10.4, 1.8, 2.6, 4.9, "Layer 4\nDeployment QA", [
    "Latency & power budget",
    "On-device memory",
    "Post-market drift monitoring",
])
add_footer(s, "Adapted for ECG/BP/SpO2/PPG/temp product pipelines")
add_notes(s, ["这是给你们团队最可执行的一页。"])
md_slide("Slide 10 - Verification framework", ["4-layer QA framework from representation to deployment", "Directly adaptable to monitoring products"])

# 11 roadmap
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_title(s, "90-Day Adoption Roadmap (For Your Team)")
add_card(s, 0.8, 1.9, 3.9, 4.7, "Phase 1 (Weeks 1-3)", [
    "Build unified biosignal data contract",
    "Define PPG/ECG/BP/SpO2/temp segmentation",
    "Set participant-level train/val/test governance",
])
add_card(s, 4.95, 1.9, 3.9, 4.7, "Phase 2 (Weeks 4-7)", [
    "Pretrain baseline encoder(s)",
    "Run linear probe matrix for core biomarkers",
    "Establish subgroup and shift dashboards",
])
add_card(s, 9.1, 1.9, 3.9, 4.7, "Phase 3 (Weeks 8-12)", [
    "Task heads + calibration + thresholds",
    "On-device profiling (latency/power)",
    "Pilot deployment + drift alert policy",
])
add_footer(s, "A practical translation of the paper into product execution")
add_notes(s, ["如果需要，我可以把这页再拆成甘特图版本。"])
md_slide("Slide 11 - Roadmap", ["90-day execution plan for adapting the framework", "Covers data, training, validation, deployment"])

# 12 risks
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, LIGHT_BG)
add_title(s, "Known Limits and What to Watch")
add_bullets(s, 0.9, 1.8, 11.8, 4.8, [
    "Self-reported labels were used for many downstream targets (label quality ceiling)",
    "Sensor-level raw data/code sharing is constrained by privacy/consent policies",
    "Embedding quality proxies (SER/InfoNCE) correlate with downstream metrics but are not sufficient",
    "Modality heterogeneity means one-size-fits-all training policy can underperform",
    "Clinical-grade use still needs independent external validation and regulatory evidence",
], size=20)
warn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(6.0), Inches(11.9), Inches(0.9))
warn.fill.solid(); warn.fill.fore_color.rgb = RGBColor(255, 244, 230); warn.line.color.rgb = ACCENT
wt = warn.text_frame
wt.text = "Recommendation: treat this paper as a strong engineering blueprint, not as a standalone clinical validation package."
wt.paragraphs[0].font.name = "Calibri"; wt.paragraphs[0].font.size = Pt(16); wt.paragraphs[0].font.bold = True; wt.paragraphs[0].font.color.rgb = RGBColor(120, 72, 19)
add_footer(s, "Section 6 + reproducibility statement")
add_notes(s, ["这页用于管理层预期管理和风险沟通。"])
md_slide("Slide 12 - Limits", ["Important limitations for productization and evidence claims", "Need external validation beyond this paper"])

# 13 conclusion
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)
add_title(s, "Final Takeaways", "How Apple’s experience can inform your monitoring AI strategy", dark=True)
add_bullets(s, 0.9, 2.0, 11.8, 3.8, [
    "Use representation learning as a shared backbone for multi-parameter monitoring",
    "Prioritize participant-level data governance and split strategy",
    "Design modality-specific training recipes (PPG != ECG != BP/SpO2/temp)",
    "Institutionalize a layered validation stack before deployment",
], color=RGBColor(234, 242, 252), size=24)
closing = s.shapes.add_textbox(Inches(0.9), Inches(6.1), Inches(11.8), Inches(0.8)).text_frame
closing.text = "Next step: run an internal pilot on your own ECG/PPG/BP datasets with this exact framework scaffold."
closing.paragraphs[0].font.name = "Calibri"
closing.paragraphs[0].font.size = Pt(17)
closing.paragraphs[0].font.color.rgb = MINT
closing.paragraphs[0].font.bold = True
add_footer(s, "Prepared from: Abbaspourazad et al., ICLR 2024")
add_notes(s, ["收尾聚焦：从单篇论文转化为你的团队执行方法。"])
md_slide("Slide 13 - Takeaways", ["Representation-first, modality-aware, validation-layered development strategy"])

# 14 appendix numbers
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_title(s, "Appendix: Key Numbers Reference")
add_card(s, 0.8, 1.8, 6.0, 4.8, "Paper constants used in this deck", [
    "Participants: PPG 141,207 | ECG 106,643",
    "Segments: PPG 19,854,101 | ECG 3,743,679",
    "Encoder params: 3.3M (PPG), 2.5M (ECG)",
    "Train split: 80/10/10 by participant",
    "InfoNCE temp 0.04, KoLeo weight 0.1",
    "Compute: 32 x A100 GPUs",
])
add_card(s, 7.1, 1.8, 5.6, 4.8, "Top quantitative outcomes", [
    "PPG age cls AUC: 0.976 (participant-level pairing)",
    "ECG age cls AUC: 0.916 (participant-level pairing)",
    "PPG sex cls AUC: 0.993 | ECG sex cls AUC: 0.951",
    "SER best: Ours > Ours(no KoLeo) > SimCLR > BYOL",
    "PPG embeddings generally outperformed baseline feature set",
])
add_footer(s, "All values from extracted paper text/table content")
add_notes(s, ["作为答疑备用页。"])
md_slide("Slide 14 - Appendix", ["Consolidated key constants and metrics from the paper"])

prs.save(OUT_PPTX)
with open(OUT_MD, 'w', encoding='utf-8') as f:
    f.write('# Large-Scale Training Deck\n\n')
    f.write('\n'.join(slides_md))
print(OUT_PPTX)
print(OUT_MD)
