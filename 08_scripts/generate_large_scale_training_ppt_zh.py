#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUT = '/Users/codex/Lab/07_reports/ppt/2024-01-29_large-scale-training-中文版.pptx'
OUT_MD = '/Users/codex/Lab/07_reports/ppt/2024-01-29_large-scale-training-中文版.slides.md'

NAVY = RGBColor(14, 36, 64)
TEAL = RGBColor(0, 137, 123)
LIGHT = RGBColor(246, 250, 252)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(26, 32, 44)
GRAY = RGBColor(90, 102, 118)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides_md = []

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def title(slide, t, s=None, dark=False):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.42), Inches(12.0), Inches(1.0)).text_frame
    p = box.paragraphs[0]
    p.text = t
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE if dark else NAVY
    if s:
        sb = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12.0), Inches(0.6)).text_frame
        sp = sb.paragraphs[0]
        sp.text = s
        sp.font.name = 'Microsoft YaHei'
        sp.font.size = Pt(16)
        sp.font.color.rgb = RGBColor(220, 230, 240) if dark else GRAY

def bullets(slide, x, y, w, h, items, size=20, color=DARK):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)

def card(slide, x, y, w, h, t, lines, fill=WHITE):
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.color.rgb = RGBColor(220, 230, 240)
    tt = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.15), Inches(w-0.4), Inches(0.45)).text_frame
    tp = tt.paragraphs[0]
    tp.text = t
    tp.font.name = 'Microsoft YaHei'
    tp.font.size = Pt(17)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    bt = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.62), Inches(w-0.4), Inches(h-0.78)).text_frame
    bt.clear()
    for i, ln in enumerate(lines):
        p = bt.paragraphs[0] if i == 0 else bt.add_paragraph()
        p.text = ln
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.space_after = Pt(3)

def footer(slide, txt):
    f = slide.shapes.add_textbox(Inches(0.7), Inches(7.08), Inches(12.0), Inches(0.28)).text_frame
    p = f.paragraphs[0]
    p.text = txt
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY

def note(slide, lines):
    ns = slide.notes_slide.notes_text_frame
    ns.clear()
    ns.text = '\n'.join(lines)

def md(h, arr):
    slides_md.append(f'## {h}\n')
    for a in arr:
        slides_md.append(f'- {a}')
    slides_md.append('')

# 1
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
title(s, '大规模可穿戴生理信号基础模型训练框架', '基于 ICLR 2024（Apple）论文的监护设备算法解读', dark=True)
bullets(s, 0.8, 2.1, 8.8, 3.0, [
    '目标：从“标签稀缺”走向“可复用算法底座”',
    '数据规模：141,207人，约20M PPG片段，约3.75M ECG片段',
    '适用：心电、PPG、血氧、血压、体温等多参数监护AI管线',
], size=22, color=RGBColor(232,240,248))
box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.8), Inches(1.8), Inches(2.8), Inches(3.8))
box.fill.solid(); box.fill.fore_color.rgb = TEAL; box.line.fill.background()
k = s.shapes.add_textbox(Inches(10.05), Inches(2.0), Inches(2.3), Inches(3.4)).text_frame
for i,t in enumerate(['141K','参与者','20M','PPG片段','3.75M','ECG片段']):
    p = k.paragraphs[0] if i==0 else k.add_paragraph()
    p.text = t
    p.font.name='Microsoft YaHei'; p.font.size=Pt(30 if i%2==0 else 14); p.font.bold=(i%2==0); p.font.color.rgb=WHITE
footer(s,'来源：Abbaspourazad et al., ICLR 2024')
md('第1页 标题', ['论文框架+监护设备落地视角'])

# 2
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT)
title(s, '这篇论文为什么值得监护设备团队重点看')
card(s,0.8,1.7,4.0,2.2,'行业痛点',['医疗标签贵且慢','自由生活场景噪声大','单任务模型开发效率低'])
card(s,5.0,1.7,4.0,2.2,'论文解法',['无监督预训练构建表征底座','一套编码器服务多任务','减少对高质量标签依赖'])
card(s,9.2,1.7,3.5,2.2,'你的收益',['缩短算法迭代周期','统一训练-验证方法','更容易做跨参数迁移'])
bullets(s,0.9,4.3,11.7,2.0,['核心思想：先学“通用生理表示”，再做任务头微调。'],size=22)
footer(s,'框架视角：数据 -> 预训练 -> 下游任务 -> 验证 -> 部署')
md('第2页 价值', ['从单任务开发转向平台化表征学习'])

# 3
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
title(s, '数据与研究设置（AHMS）')
card(s,0.8,1.8,6.1,4.9,'数据规模',['PPG：141,207人，19,854,101片段','ECG：106,643人，3,743,679片段','时间跨度：PPG 890天，ECG 1240天','切分：按参与者80/10/10，无重叠'])
card(s,7.1,1.8,5.5,4.9,'采集与预处理',['PPG：60秒，64/256Hz，4通道','ECG：30秒，Lead-I，512Hz->128Hz','PPG预处理：暗电流扣除+带通+标准化','ECG预处理：内部工具+标准化'])
footer(s,'这部分对应论文第4节与Table 1')
md('第3页 数据设计', ['大规模纵向可穿戴数据是方法成立前提'])

# 4
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT)
title(s, '算法总流程（系统视角）')
card(s,0.8,2.0,2.3,3.8,'1 数据构建',['分模态切片','参与者均衡采样'])
card(s,3.4,2.0,2.3,3.8,'2 数据增强',['cutout/warp/noise','PPG与ECG不同策略'])
card(s,6.0,2.0,2.3,3.8,'3 编码器',['1D-EfficientNet','256维嵌入'])
card(s,8.6,2.0,2.3,3.8,'4 训练目标',['InfoNCE + KoLeo','L2归一化'])
card(s,11.2,2.0,1.8,3.8,'5 动量分支',['EMA更新','稳定训练'])
for x in [3.1,5.7,8.3,10.9]:
    c = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(3.55), Inches(0.25), Inches(0.45))
    c.fill.solid(); c.fill.fore_color.rgb = TEAL; c.line.fill.background()
footer(s,'可直接映射为你们算法架构图')
md('第4页 系统流程', ['5阶段SSL流程，可复用到监护参数算法开发'])

# 5
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
title(s, '训练目标与关键超参数')
card(s,0.8,1.8,6.2,2.3,'目标函数',['L = InfoNCE + KoLeo正则','温度τ=0.04，KoLeo权重=0.1','优化器Adam，32×A100训练'])
card(s,0.8,4.3,6.2,2.2,'设计意义',['InfoNCE：拉近正样本、拉远负样本','KoLeo：提升嵌入分布熵，减少塌陷风险','动量分支：提升训练稳定性'])
card(s,7.3,1.8,5.2,4.7,'增强策略（论文原值）',['PPG：cutout0.4 / mag-warp0.25 / noise0.25 / channel permute0.25 / time-warp0.15','ECG：cutout0.8 / mag-warp0.5 / noise0.5 / time-warp0.3','结论：ECG增强更强，因为其个体内变化更小'])
footer(s,'第3节 + 实现细节')
md('第5页 训练策略', ['InfoNCE+KoLeo+动量分支是核心组合'])

# 6
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT)
title(s, '关键结果：参与者级正样本对明显优于片段级')
card(s,0.8,1.8,5.9,4.9,'PPG（参与者级）',['年龄分类AUC 0.976（片段级0.900）','年龄回归MAE 3.19（片段级6.60）','BMI分类AUC 0.918（片段级0.766）','BMI回归MAE 2.54（片段级4.05）','性别分类AUC 0.993（片段级0.870）'])
card(s,7.0,1.8,5.9,4.9,'ECG（参与者级）',['年龄分类AUC 0.916（片段级0.783）','年龄回归MAE 6.33（片段级9.01）','BMI分类AUC 0.797（片段级0.734）','BMI回归MAE 3.72（片段级4.11）','性别分类AUC 0.951（片段级0.854）'])
footer(s,'直接来自Table 2')
md('第6页 结果', ['正样本对构造策略是高杠杆设计点'])

# 7
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
title(s, 'PPG 与 ECG 的方法学差异（非常关键）')
card(s,0.8,1.8,4.1,4.8,'论文观察',['ECG在t-SNE聚类更紧','ECG InfoNCE更低（预训练更“容易”）','ECG离散度比PPG更低','但PPG对多种健康目标预测更强'])
card(s,5.2,1.8,3.8,4.8,'解释',['不同模态信息结构不同','同一训练配方并非最优','需要模态特异设计'])
card(s,9.3,1.8,3.8,4.8,'对你们的启示',['PPG/ECG/BP/SpO2/体温应分模态策略','各模态独立看训练曲线与验证面板','避免“一套参数跑所有信号”'])
footer(s,'第5.2.2节')
md('第7页 模态差异', ['多参数监护算法必须模态化设计'])

# 8
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT)
title(s, '消融：预训练框架对比（SER）')
card(s,0.8,1.8,6.0,4.8,'PPG/ECG Smooth Effective Rank',['Ours：104.13 / 113.82','Ours(去KoLeo)：101.29 / 108.84','SimCLR变体：98.87 / 103.65','BYOL变体：51.18 / 63.67'])
card(s,7.1,1.8,5.6,4.8,'结论',['本文框架整体优于SimCLR/BYOL变体','KoLeo正则确实贡献性能','无监督代理指标有价值，但不能替代临床终点'])
footer(s,'来自Table 3')
md('第8页 消融-框架', ['KoLeo正则和框架组合有效'])

# 9
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
title(s, '消融：编码器架构与部署成本')
card(s,0.8,1.8,4.0,4.8,'1D-EfficientNet',['SER：PPG 104.13 / ECG 113.82','参数量：3.3M（PPG）/2.5M（ECG）','性能-资源比优'])
card(s,5.0,1.8,4.0,4.8,'1D-ResNet',['SER：95.81 / 111.92','参数量：16.9M','更大但收益有限'])
card(s,9.2,1.8,3.9,4.8,'1D-ViT',['SER：104.89 / 114.62','参数量：7.2M','性能强但部署成本更高'])
footer(s,'来自Table 4')
md('第9页 消融-架构', ['设备端先追求性能/功耗/内存平衡'])

# 10
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT)
title(s, '可直接复用的监护算法验证大框架')
card(s,0.8,1.8,3.0,4.9,'层1 表征质量',['SER/嵌入漂移','模态训练损失','个体内外离散度'])
card(s,4.0,1.8,3.0,4.9,'层2 任务质量',['核心指标+校准','亚组表现（年龄/性别/BMI）','阈值敏感性'])
card(s,7.2,1.8,3.0,4.9,'层3 稳健性',['分布偏移切片','标签缺失/偏移校正','传感器质量分层'])
card(s,10.4,1.8,2.6,4.9,'层4 部署质量',['时延/功耗','端侧内存','上线后漂移监控'])
footer(s,'这页建议直接纳入你们研发流程文档')
md('第10页 验证框架', ['从表征到部署的四层验证闭环'])

# 11
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, WHITE)
title(s, '给你的团队：90天落地计划')
card(s,0.8,1.9,3.9,4.7,'第1阶段（1-3周）',['统一数据契约与切片规范','定义模态：ECG/PPG/BP/SpO2/体温','建立参与者级切分治理'])
card(s,5.0,1.9,3.9,4.7,'第2阶段（4-7周）',['训练基础编码器','建立线性探针评估矩阵','接入亚组与偏移监控'])
card(s,9.2,1.9,3.9,4.7,'第3阶段（8-12周）',['接任务头与阈值策略','端侧性能画像（功耗/时延）','小规模试点与漂移告警'])
footer(s,'执行导向：先小闭环，再扩展多参数')
md('第11页 落地路线', ['90天从方法验证走到试点部署'])

# 12
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, NAVY)
title(s, '结论：苹果经验可借鉴的三件事', '面向监护设备AI研发的可执行建议', dark=True)
bullets(s,0.9,2.0,11.8,3.8,[
    '先做“表征底座”，再做参数任务头，不要每个参数都从零训练',
    '训练策略必须分模态定制，尤其是PPG与ECG',
    '把验证体系前置：亚组、公平性、偏移、部署约束同时纳入',
],size=25,color=RGBColor(235,242,252))
footer(s,'Prepared for monitoring algorithm R&D')
md('第12页 总结', ['表征优先、模态化训练、验证前置'])

prs.save(OUT)
with open(OUT_MD,'w',encoding='utf-8') as f:
    f.write('# Large-Scale Training 中文版\n\n')
    f.write('\n'.join(slides_md))

print(OUT)
print(OUT_MD)
